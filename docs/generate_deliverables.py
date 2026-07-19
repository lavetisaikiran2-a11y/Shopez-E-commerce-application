import os
import re
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from pptx import Presentation

# Set up ReportLab Styles
styles = getSampleStyleSheet()

title_style = ParagraphStyle(
    'DocTitle',
    parent=styles['Normal'],
    fontName='Helvetica-Bold',
    fontSize=20,
    leading=24,
    textColor=colors.HexColor('#2874f0'),
    spaceAfter=15
)
h1_style = ParagraphStyle(
    'DocH1',
    parent=styles['Normal'],
    fontName='Helvetica-Bold',
    fontSize=14,
    leading=18,
    textColor=colors.HexColor('#333333'),
    spaceBefore=12,
    spaceAfter=6
)
h2_style = ParagraphStyle(
    'DocH2',
    parent=styles['Normal'],
    fontName='Helvetica-Bold',
    fontSize=11,
    leading=14,
    textColor=colors.HexColor('#555555'),
    spaceBefore=10,
    spaceAfter=4
)
body_style = ParagraphStyle(
    'DocBody',
    parent=styles['Normal'],
    fontName='Helvetica',
    fontSize=9,
    leading=13,
    textColor=colors.HexColor('#222222'),
    spaceAfter=5
)
bullet_style = ParagraphStyle(
    'DocBullet',
    parent=body_style,
    leftIndent=15,
    firstLineIndent=-10,
    spaceAfter=3
)
code_style = ParagraphStyle(
    'DocCode',
    parent=styles['Normal'],
    fontName='Courier',
    fontSize=8,
    leading=10,
    textColor=colors.HexColor('#a71d5d'),
    backColor=colors.HexColor('#f6f8fa'),
    borderColor=colors.HexColor('#dddddd'),
    borderWidth=0.5,
    borderPadding=4,
    spaceAfter=5
)

def parse_markdown_table(table_lines):
    rows = []
    for line in table_lines:
        if '---|' in line or '--:|' in line or '|---|' in line:
            continue
        cells = [c.strip() for c in line.split('|')[1:-1]]
        if cells:
            rows.append(cells)
    
    data = []
    for i, r in enumerate(rows):
        row_data = []
        for cell in r:
            # Strip simple markdown bold/italic tags for PDF compatibility
            clean_cell = re.sub(r'\*\*|__', '', cell)
            # Escape HTML characters to avoid paraparser crashes
            clean_cell = clean_cell.replace('&', '&amp;')
            clean_cell = clean_cell.replace('<br>', '__BR__')
            clean_cell = clean_cell.replace('<br/>', '__BR__')
            clean_cell = clean_cell.replace('<', '&lt;')
            clean_cell = clean_cell.replace('>', '&gt;')
            clean_cell = clean_cell.replace('__BR__', '<br/>')
            style = h2_style if i == 0 else body_style
            row_data.append(Paragraph(clean_cell, style))
        data.append(row_data)
        
    t = Table(data, hAlign='LEFT')
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#f2f2f2')),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('TOPPADDING', (0,0), (-1,-1), 4),
        ('LEFTPADDING', (0,0), (-1,-1), 4),
        ('RIGHTPADDING', (0,0), (-1,-1), 4),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cccccc')),
    ]))
    return t

def md_to_pdf(md_path, pdf_path):
    print(f"Generating PDF: {pdf_path} from {md_path}")
    doc = SimpleDocTemplate(pdf_path, pagesize=letter, rightMargin=45, leftMargin=45, topMargin=45, bottomMargin=45)
    story = []
    
    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()
        
    in_code = False
    code_block = []
    in_table = False
    table_lines = []
    
    for line in lines:
        stripped = line.strip()
        
        # Code Block Toggle
        if stripped.startswith('```'):
            if in_code:
                in_code = False
                code_text = '\n'.join(code_block)
                story.append(Paragraph(code_text.replace('\n', '<br/>'), code_style))
                story.append(Spacer(1, 4))
                code_block = []
            else:
                in_code = True
            continue
            
        if in_code:
            code_block.append(line.rstrip())
            continue
            
        # Table Detection
        if stripped.startswith('|'):
            in_table = True
            table_lines.append(stripped)
            continue
        elif in_table:
            in_table = False
            if table_lines:
                story.append(parse_markdown_table(table_lines))
                story.append(Spacer(1, 6))
            table_lines = []
            
        # Headings & Paragraph elements
        if stripped.startswith('# '):
            clean_text = stripped[2:].replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            story.append(Paragraph(clean_text, title_style))
            story.append(Spacer(1, 8))
        elif stripped.startswith('## '):
            clean_text = stripped[3:].replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            story.append(Paragraph(clean_text, h1_style))
            story.append(Spacer(1, 6))
        elif stripped.startswith('### '):
            clean_text = stripped[4:].replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            story.append(Paragraph(clean_text, h2_style))
            story.append(Spacer(1, 4))
        elif stripped.startswith('- ') or stripped.startswith('* '):
            # Strip simple markdown links from bullet points
            bullet_text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', stripped[2:])
            clean_text = bullet_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            story.append(Paragraph('&bull; ' + clean_text, bullet_style))
        elif stripped == '':
            story.append(Spacer(1, 3))
        else:
            # Strip simple markdown links and brackets
            clean_text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', stripped)
            clean_text = clean_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            story.append(Paragraph(clean_text, body_style))
            
    doc.build(story)
    print(f"Successfully generated {pdf_path}")

def md_to_pptx(md_path, pptx_path):
    print(f"Generating PPTX: {pptx_path} from {md_path}")
    prs = Presentation()
    
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
        
    slides_raw = content.split('---')
    
    for slide_data in slides_raw:
        lines = [l.strip() for l in slide_data.strip().split('\n') if l.strip()]
        if not lines:
            continue
            
        title_text = ""
        body_text_lines = []
        
        for line in lines:
            if line.startswith('# ') or line.startswith('## '):
                title_text = re.sub(r'^[#\s]+', '', line)
            elif line.startswith('- ') or line.startswith('* '):
                body_text_lines.append(line)
            elif not line.startswith('>') and not line.startswith('Visuals:'):
                body_text_lines.append(line)
                
        if "Overview" in title_text or "Presentation" in title_text:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title_text
            if body_text_lines:
                slide.placeholders[1].text = '\n'.join(body_text_lines)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title_text or "ShopEZ Highlights"
            if body_text_lines:
                slide.placeholders[1].text = '\n'.join(body_text_lines)
                
    prs.save(pptx_path)
    print(f"Successfully generated {pptx_path}")

if __name__ == '__main__':
    docs_dir = os.path.dirname(os.path.abspath(__file__))
    
    # Generate Reports
    md_to_pdf(os.path.join(docs_dir, "Final_Project_Report.md"), os.path.join(docs_dir, "Final_Project_Report.pdf"))
    md_to_pdf(os.path.join(docs_dir, "Testing_Report.md"), os.path.join(docs_dir, "Testing_Report.pdf"))
    md_to_pdf(os.path.join(docs_dir, "UAT_Report.md"), os.path.join(docs_dir, "UAT_Report.pdf"))
    
    # Generate presentation
    md_to_pptx(os.path.join(docs_dir, "ShopEZ_Presentation.md"), os.path.join(docs_dir, "ShopEZ_Presentation.pptx"))
